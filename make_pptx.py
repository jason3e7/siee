#!/usr/bin/env python3
"""Generate siee.pptx — SIEE project presentation."""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0f, 0x17, 0x2a)
BG2     = RGBColor(0x1e, 0x29, 0x3b)
BG3     = RGBColor(0x09, 0x11, 0x1e)
ACCENT  = RGBColor(0x4a, 0xde, 0x80)
DANGER  = RGBColor(0xf8, 0x71, 0x71)
WARN    = RGBColor(0xfb, 0xbf, 0x24)
BLUE    = RGBColor(0x60, 0xa5, 0xfa)
PURPLE  = RGBColor(0xc0, 0x84, 0xfc)
MUTED   = RGBColor(0x94, 0xa3, 0xb8)
WHITE   = RGBColor(0xf1, 0xf5, 0xf9)
BORDER  = RGBColor(0x33, 0x41, 0x55)
CODE_FG = RGBColor(0xa5, 0xf3, 0xfc)

# ── Layout ────────────────────────────────────────────────────────────────────
W        = Inches(13.33)
H        = Inches(7.5)
ML       = Inches(0.65)
MT       = Inches(0.45)
CW       = Inches(12.03)
HW       = Inches(5.86)
GAP      = Inches(0.31)
R        = ML + HW + GAP
BODY_TOP = Inches(1.18)
BODY_H   = Inches(5.9)
CODE_PT  = Pt(10.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def _set_txbody_margins(txBody, l=Pt(8), r=Pt(8), t=Pt(6), b=Pt(6), anchor='t'):
    bp = txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('lIns', str(int(l)))
        bp.set('rIns', str(int(r)))
        bp.set('tIns', str(int(t)))
        bp.set('bIns', str(int(b)))
        bp.set('anchor', anchor)


def _set_anchor(txBody, anchor='t'):
    bp = txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)


def txtbox(slide, x, y, w, h, text, size=Pt(14), color=WHITE,
           bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Noto Sans TC"
    return tb


def heading(slide, text, y=MT):
    txtbox(slide, ML, y, CW, Inches(0.62), text,
           size=Pt(30), color=WHITE, bold=True)
    rule = slide.shapes.add_shape(1, ML, y + Inches(0.63), CW, Pt(2.5))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def rect(slide, x, y, w, h, fill_color=BG2, border_color=BORDER, border_pt=1.2):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def card(slide, x, y, w, h, title="", body="",
         border=BORDER, title_color=ACCENT, body_color=WHITE,
         body_size=Pt(13.5), fill=BG2):
    shape = rect(slide, x, y, w, h, fill_color=fill, border_color=border)
    tf = shape.text_frame
    tf.word_wrap = True
    _set_txbody_margins(tf._txBody, l=Pt(11), r=Pt(11), t=Pt(9), b=Pt(9))
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Noto Sans TC"
    if body:
        p2 = tf.paragraphs[0] if not title else tf.add_paragraph()
        if title:
            p2.space_before = Pt(4)
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = body
        run2.font.size = body_size
        run2.font.color.rgb = body_color
        run2.font.name = "Noto Sans TC"
    return shape


def code_block(slide, x, y, w, h, code: str, label="", label_color=MUTED):
    if label:
        txtbox(slide, x, y, w, Inches(0.3), label,
               size=Pt(12), color=label_color)
        y += Inches(0.3); h -= Inches(0.3)
    shape = rect(slide, x, y, w, h, fill_color=BG3, border_color=BORDER, border_pt=1)
    tf = shape.text_frame
    tf.word_wrap = False
    _set_txbody_margins(tf._txBody, l=Pt(11), r=Pt(11), t=Pt(9), b=Pt(9))
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(17)
        run = p.add_run()
        run.text = line
        run.font.size = CODE_PT
        run.font.color.rgb = CODE_FG
        run.font.name = "Courier New"
    return shape


def table(slide, x, y, w, col_fracs: list,
          headers: list, rows: list,
          row_h=Inches(0.44), cell_size=Pt(13)):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, row_h * n_rows)
    total_fracs = sum(col_fracs)
    for ci, frac in enumerate(col_fracs):
        tbl.table.columns[ci].width = int(w * frac / total_fracs)

    def fmt(cell, text, bold=False, fg=WHITE, bg=BG2, size=cell_size, align=PP_ALIGN.LEFT):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        tf = cell.text_frame; tf.word_wrap = True
        _set_txbody_margins(tf._txBody, l=Pt(7), r=Pt(5), t=Pt(5), b=Pt(5))
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = size; run.font.bold = bold
        run.font.color.rgb = fg; run.font.name = "Noto Sans TC"

    HDR_BG = RGBColor(0x0d, 0x24, 0x18)
    for ci, h in enumerate(headers):
        fmt(tbl.table.cell(0, ci), h, bold=True, fg=ACCENT, bg=HDR_BG)
    for ri, row in enumerate(rows):
        bg = RGBColor(0x14, 0x1e, 0x32) if ri % 2 == 0 else BG
        for ci, val in enumerate(row):
            fg = RGBColor(0x6e, 0xe7, 0xb7) if ci == 0 else WHITE
            fmt(tbl.table.cell(ri + 1, ci), str(val), fg=fg, bg=bg)
    return tbl


def stat_box(slide, x, y, w, h, number, label):
    shape = rect(slide, x, y, w, h, fill_color=BG2, border_color=BORDER)
    tf = shape.text_frame; tf.word_wrap = True
    _set_anchor(tf._txBody, 'ctr')
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number; r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = ACCENT; r.font.name = "Noto Sans TC"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label; r2.font.size = Pt(13)
    r2.font.color.rgb = MUTED; r2.font.name = "Noto Sans TC"


def flow_row(slide, items: list[tuple], x, y, h=Inches(0.5)):
    cx = x
    for label, is_cmd, color in items:
        is_arr = label in ("→", "←")
        fw = Inches(0.3) if is_arr else Inches(1.4)
        if is_arr:
            shape = rect(slide, cx, y, fw, h, fill_color=BG, border_color=None)
        else:
            bg = RGBColor(0x0e, 0x2a, 0x1a) if is_cmd else BG2
            bc = ACCENT if is_cmd else BORDER
            shape = rect(slide, cx, y, fw, h, fill_color=bg, border_color=bc)
        tf = shape.text_frame
        _set_anchor(tf._txBody, 'ctr')
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(12.5)
        run.font.bold = is_cmd
        run.font.color.rgb = color if not is_arr else MUTED
        run.font.name = "Noto Sans TC"
        cx += fw + Inches(0.055)


# ══════════════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════════════

def s01_cover(prs):
    slide = blank(prs)
    START = Inches((7.5 - 4.6) / 2)

    txtbox(slide, ML, START, CW, Inches(0.32),
           "Open Source  ·  Python  ·  Flask  ·  MCP",
           size=Pt(13.5), color=MUTED, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, START + Inches(0.38), CW, Inches(1.05),
           "SIEE",
           size=Pt(62), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, START + Inches(1.5), CW, Inches(0.52),
           "Secret Isolation Execution Environment",
           size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, START + Inches(2.08), CW, Inches(0.45),
           "讓 AI agent 可以執行真實 API 測試，但永遠摸不到 secret 的值。",
           size=Pt(14.5), color=MUTED, align=PP_ALIGN.CENTER)

    sw = Inches(2.75); sh = Inches(1.6)
    sx = ML + Inches(0.26)
    sy = START + Inches(2.75)
    for num, lbl in [("3", "MCP Tools"), ("23", "自動化測試"), ("2", "獨立 Server"), ("0", "洩漏的 Secret")]:
        stat_box(slide, sx, sy, sw, sh, num, lbl)
        sx += sw + Inches(0.29)


def s02_problem(prs):
    slide = blank(prs)
    heading(slide, "問題：AI agent 需要 token，但不能讓它看到")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.32),
           "用 AI 開發需要整合外部 API 的專案時，一定會遇到這個矛盾……",
           size=Pt(14), color=MUTED)

    cw = Inches(3.9); cy = BODY_TOP + Inches(0.45); ch = Inches(2.1)
    cx = ML
    for title, bc, body in [
        ("🔑 Token 不能進對話",   DANGER,
         "把 API key 貼給 AI\n等於讓它進了你的 context\n一旦洩漏無法收回"),
        ("🧪 測試需要真實 token", WARN,
         "Mock 測試通過≠整合正確\n真實 API 才能驗證行為\n假環境掩蓋真正的錯誤"),
        ("🤖 AI 寫的 code 要跑起來", BLUE,
         "AI 寫完 code 要立刻驗證\n但執行環境有 secret\nAI 不該直接碰執行環境"),
    ]:
        card(slide, cx, cy, cw, ch, title=title, body=body,
             border=bc, title_color=bc, body_size=Pt(14))
        cx += cw + Inches(0.17)

    card(slide, ML, cy + ch + Inches(0.3), CW, Inches(0.95),
         body="💡  不用貼來貼：AI 自動部署 code、觸發執行、拿回結果，全流程自動化——secret 值永遠不進 AI 的 context。",
         border=ACCENT, body_color=ACCENT, body_size=Pt(15),
         fill=RGBColor(0x0a, 0x1f, 0x12))


def s03_solution(prs):
    slide = blank(prs)
    heading(slide, "解法：SIEE 作為隔離執行層")

    flow_items = [
        ("AI Agent",    False, WHITE),  ("→", False, MUTED),
        ("deploy",      True,  ACCENT), ("→", False, MUTED),
        ("SIEE Server", False, BLUE),   ("→", False, MUTED),
        ("workspace/",  False, MUTED),
    ]
    flow_row(slide, flow_items, ML, BODY_TOP + Inches(0.1))

    flow_items2 = [
        ("AI Agent",    False, WHITE),  ("→", False, MUTED),
        ("exec",        True,  ACCENT), ("→", False, MUTED),
        ("SIEE Server", False, BLUE),   ("→", False, MUTED),
        ("subprocess",  False, PURPLE), ("→", False, MUTED),
        ("secret 注入",  False, WARN),
    ]
    flow_row(slide, flow_items2, ML, BODY_TOP + Inches(0.75))

    flow_items3 = [
        ("AI Agent",    False, WHITE),  ("→", False, MUTED),
        ("get_log",     True,  ACCENT), ("←", False, MUTED),
        ("SIEE Server", False, BLUE),   ("←", False, MUTED),
        ("stdout/stderr", False, MUTED),
    ]
    flow_row(slide, flow_items3, ML, BODY_TOP + Inches(1.4))

    txtbox(slide, ML, BODY_TOP + Inches(2.1), CW, Inches(0.3),
           "AI 只看到執行結果，secret 的值永遠不出現在 context window 裡。",
           size=Pt(13.5), color=MUTED)

    cw = Inches(5.86); ch = Inches(2.4)
    cy = BODY_TOP + Inches(2.55)

    card(slide, ML, cy, cw, ch,
         title="SIEE Server（這台機器）",
         body="• secret 以 env var 形式存在\n• subprocess 執行時注入\n• AI 從未知道值是什麼\n• log 寫到 logs/ 資料夾",
         border=ACCENT, title_color=ACCENT, body_size=Pt(13.5))

    card(slide, R, cy, cw, ch,
         title="AI Agent（另一台機器）",
         body="• 只知道 SIEE 的 IP + port\n• 透過 MCP tools 操作\n• 拿到的只有 stdout/stderr\n• 無法存取 env var",
         border=BLUE, title_color=BLUE, body_size=Pt(13.5))


def s04_philosophy(prs):
    slide = blank(prs)
    heading(slide, "設計哲學：和 GitHub Actions Secrets 一樣")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.35),
           "這不是新概念——CI/CD 世界已經用這個模式幾十年了。",
           size=Pt(14), color=MUTED)

    headers = ["", "GitHub Actions Secrets", "SIEE"]
    rows = [
        ["誰寫 code",       "開發者",                     "AI Agent"],
        ["Secret 存在哪",   "GitHub repository settings", "SIEE server env var"],
        ["如何注入",         "Runner 執行時注入環境變數",   "subprocess 執行時注入"],
        ["作者看到什麼",     "workflow 結果（pass/fail）", "stdout / stderr"],
        ["作者能讀 secret？", "❌ 不行",                   "❌ 不行"],
    ]
    table(slide, ML, BODY_TOP + Inches(0.45), CW,
          [1.5, 3.0, 3.0], headers, rows, row_h=Inches(0.5))

    card(slide, ML, BODY_TOP + Inches(3.45), CW, Inches(0.95),
         body="「Code 可以『使用』secret，但作者（或 AI）永遠不會『看到』secret 的值。」",
         border=ACCENT, body_color=WHITE, body_size=Pt(16),
         fill=RGBColor(0x0a, 0x1f, 0x12))


def s05_threat_model(prs):
    slide = blank(prs)
    heading(slide, "威脅模型：防什麼，不防什麼")

    cw = Inches(5.86); ch = Inches(4.5)
    cy = BODY_TOP + Inches(0.1)

    card(slide, ML, cy, cw, ch,
         title="✅  SIEE 能防的",
         body=("• .env 檔案被 AI 讀進 context\n"
               "• config 檔含 token 被當作 context 讀取\n"
               "• secret 被貼進對話裡\n"
               "• prompt injection 讓 AI 說出 token\n"
               "  （因為 AI 根本不知道值）\n"
               "• token 出現在 AI 的 context window"),
         border=ACCENT, title_color=ACCENT, body_size=Pt(14))

    card(slide, R, cy, cw, ch,
         title="❌  SIEE 擋不住的",
         body=("• AI 故意寫出 print(os.environ['TOKEN'])\n"
               "  → log 裡會出現 token 值\n\n"
               "• 安全假設：\n"
               "  AI agent 是「無意識的」\n"
               "  不會主動想要竊取 token\n"
               "  只需防「被動洩漏」"),
         border=DANGER, title_color=DANGER, body_size=Pt(14))

    txtbox(slide, ML, cy + ch + Inches(0.2), CW, Inches(0.3),
           "SIEE 是 passive isolation，不是 adversarial sandbox。",
           size=Pt(13), color=MUTED)


def s06_api(prs):
    slide = blank(prs)
    heading(slide, "REST API：三個端點")

    ch = Inches(1.55); cy = BODY_TOP + Inches(0.1); cw = CW

    card(slide, ML, cy, cw, ch,
         title="POST  /deploy",
         body="上傳檔案（multipart），整個替換 workspace/。同一個專案，每次 deploy 都是全新覆蓋。\n回傳：{\"status\": \"deployed\", \"files\": [...]}",
         border=BLUE, title_color=BLUE, body_size=Pt(13.5))

    cy += ch + Inches(0.18)
    card(slide, ML, cy, cw, ch,
         title="POST  /exec",
         body="執行白名單內的指令，非同步回傳 exec_id。Body: {\"command\": \"pytest\", \"args\": [\"-v\"]}。\n回傳：{\"exec_id\": \"...\", \"status\": \"RUNNING\"}",
         border=ACCENT, title_color=ACCENT, body_size=Pt(13.5))

    cy += ch + Inches(0.18)
    card(slide, ML, cy, cw, ch,
         title="GET  /logs/{exec_id}",
         body="取得執行狀態與 stdout/stderr。狀態為 RUNNING / DONE / ERROR，輪詢直到非 RUNNING。\n回傳：{\"exec_id\": \"...\", \"status\": \"DONE\", \"log\": \"...\"}",
         border=PURPLE, title_color=PURPLE, body_size=Pt(13.5))

    txtbox(slide, ML, Inches(6.9), CW, Inches(0.35),
           "※ MVP 展示版本：核心隔離機制已完整，指令白名單與存取控制可依需求擴充。",
           size=Pt(12), color=MUTED, italic=True)


def s07_allowed_commands(prs):
    slide = blank(prs)
    heading(slide, "安全白名單：ALLOWED_COMMANDS")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.32),
           "AI 只能呼叫白名單內的指令，直接防止任意 shell 指令注入。",
           size=Pt(14), color=MUTED)

    code = ("# server.py 頂端，server owner 自行設定\n"
            "ALLOWED_COMMANDS = {\n"
            "    \"pytest\": [sys.executable, \"-m\", \"pytest\"],\n"
            "    \"run\":    [sys.executable, \"main.py\"],\n"
            "}\n\n"
            "# 呼叫方式\n"
            "POST /exec\n"
            "{\"command\": \"pytest\", \"args\": [\"-v\", \"tests/\"]}\n\n"
            "# 不在白名單 → 400 Bad Request\n"
            "{\"command\": \"rm\", \"args\": [\"-rf\", \"/\"]}\n"
            "→ {\"error\": \"command not allowed\", \"available\": [\"pytest\", \"run\"]}")

    code_block(slide, ML, BODY_TOP + Inches(0.4), CW, Inches(3.9), code)

    cw = Inches(3.75)
    cy = BODY_TOP + Inches(4.45)
    cx = ML
    for title, bc, body in [
        ("指令白名單",   ACCENT, "只有 server owner 定義的指令可以執行"),
        ("args 可選",    BLUE,   "額外參數 append 在 hardcoded 指令後"),
        ("sys.executable", PURPLE, "使用 venv 的 Python，套件保持一致"),
    ]:
        card(slide, cx, cy, cw, Inches(0.9),
             title=title, body=body, border=bc, title_color=bc, body_size=Pt(13))
        cx += cw + Inches(0.14)


def s08_mcp(prs):
    slide = blank(prs)
    heading(slide, "MCP 整合：AI Agent 零設定接入")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.32),
           "mcp_server.py 透過 SSE transport 暴露三個 MCP tools，AI agent 機器只需要加一段設定。",
           size=Pt(14), color=MUTED)

    cfg = ("# ~/.claude/settings.json  或  .claude/settings.json\n"
           "{\n"
           "  \"mcpServers\": {\n"
           "    \"siee\": {\n"
           "      \"type\": \"sse\",\n"
           "      \"url\":  \"http://192.168.0.2:5001/sse\"\n"
           "    }\n"
           "  }\n"
           "}")

    code_block(slide, ML, BODY_TOP + Inches(0.42), HW, Inches(3.3), cfg,
               "AI Agent 機器設定", label_color=ACCENT)

    tools_code = ("# AI agent 可用的 MCP tools\n\n"
                  "deploy(files={\"main.py\": \"...\"})\n"
                  "# → 上傳 code 到 SIEE\n\n"
                  "exec_command(command=\"pytest\",\n"
                  "             args=[\"-v\"])\n"
                  "# → 回傳 exec_id\n\n"
                  "get_log(exec_id=\"...\")\n"
                  "# → 輪詢直到 DONE/ERROR")

    code_block(slide, R, BODY_TOP + Inches(0.42), HW, Inches(3.3), tools_code,
               "AI 使用方式", label_color=BLUE)

    cw = Inches(5.86)
    cy = BODY_TOP + Inches(3.92)
    card(slide, ML, cy, cw, Inches(1.3),
         title="SIEE Server（這台機器）",
         body="python server.py      # port 5000\npython mcp_server.py  # port 5001\n環境變數：SIEE_URL、MCP_PORT 可覆寫",
         border=ACCENT, title_color=ACCENT, body_size=Pt(13.5))

    card(slide, R, cy, cw, Inches(1.3),
         title="SSE 是什麼",
         body="Server-Sent Events：單向 HTTP 長連線\nMCP 用它讓 Claude 保持與 tool server 的連線\n不需要 WebSocket，標準 HTTP 即可",
         border=BLUE, title_color=BLUE, body_size=Pt(13.5))


def s09_case_env(prs):
    slide = blank(prs)
    heading(slide, "實戰案例：AI 探查 SIEE 環境")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.32),
           "AI 遇到 ModuleNotFoundError 後，部署探針腳本自行偵測執行環境，直到找出問題根因。",
           size=Pt(14), color=MUTED)

    probe_code = ("# AI 部署的環境探針\n"
                  "import sys, subprocess\n"
                  "print('python:', sys.executable)\n"
                  "r = subprocess.run(\n"
                  "    [sys.executable, '-m', 'pip', 'list'],\n"
                  "    capture_output=True, text=True)\n"
                  "print(r.stdout)\n"
                  "try:\n"
                  "    import requests\n"
                  "    print('requests OK:', requests.__version__)\n"
                  "except ImportError as e:\n"
                  "    print('requests FAIL:', e)\n"
                  "    print('sys.path:', sys.path)")

    code_block(slide, ML, BODY_TOP + Inches(0.42), HW, Inches(3.7), probe_code,
               "AI 部署的探針腳本", label_color=BLUE)

    stdout_code = ("# Probe 1：系統套件，無 requests\n"
                   "python: /home/user/siee/venv/bin/python\n"
                   "版本: 3.12.3, 80+ 系統套件\n"
                   "requests  ← 不在列表！\n\n"
                   "# Probe 2：路徑確認\n"
                   "requests FAIL: No module named 'requests'\n"
                   "sys.path: [..., venv/site-packages]\n\n"
                   "# Probe 3：修復後確認\n"
                   "python: /home/user/siee/venv/bin/python\n"
                   "requests OK: 2.34.2")

    code_block(slide, R, BODY_TOP + Inches(0.42), HW, Inches(3.7), stdout_code,
               "SIEE 回傳的 stdout（AI 看到的）", label_color=ACCENT)

    cw = Inches(3.75); cy = BODY_TOP + Inches(4.35)
    cx = ML
    for title, bc, body in [
        ("5 次 ModuleNotFoundError", DANGER,
         "AI 不斷重試 import requests\n才意識到環境有異"),
        ("部署探針自我偵測",          WARN,
         "AI 寫探針 code，透過 SIEE\n取得 python 路徑與套件清單"),
        ("venv 缺少套件 → 修復",     ACCENT,
         "發現 venv 未安裝 requests\n修復後 requests OK: 2.34.2"),
    ]:
        card(slide, cx, cy, cw, Inches(1.0),
             title=title, body=body, border=bc, title_color=bc, body_size=Pt(13))
        cx += cw + Inches(0.14)


def s10_case_censys(prs):
    slide = blank(prs)
    heading(slide, "實戰案例：Censys API v3 整合開發")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.32),
           "AI 透過 SIEE 開發 Censys /v3/global/search/query 整合，API key 全程不出現在 context。",
           size=Pt(14), color=MUTED)

    cw4 = Inches(2.85); cy4 = BODY_TOP + Inches(0.42)
    cx4 = ML
    for num, lbl in [("594", "Total Hits"), ("6", "分頁數"), ("2", "資料型態"), ("3", "Retry 次數")]:
        stat_box(slide, cx4, cy4, cw4, Inches(1.15), num, lbl)
        cx4 += cw4 + Inches(0.17)

    query_code = ("# IP 查詢（geolocation + ASN）\n"
                  "host.ip=\"8.8.8.8\" or\n"
                  "host.ip=\"8.8.4.4\" or host.ip=\"1.1.1.1\"\n"
                  "→ 3 hits：Google DNS、位置、ASN\n\n"
                  "# 憑證搜尋（webproperty）\n"
                  "webproperty.ip=\"93.184.216.34\"\n"
                  "→ 594 hits / 6 pages\n"
                  "→ next_page_token 分頁\n"
                  "→ Resuming session 斷點續傳\n"
                  "20260522t021403_result.json")

    code_block(slide, ML, BODY_TOP + Inches(1.72), HW, Inches(3.8), query_code,
               "實際查詢語法與結果", label_color=ACCENT)

    challenges_code = ("# 挑戰 1：ModuleNotFoundError\n"
                       "import requests  # FAIL × 5\n"
                       "→ 探查環境後修復\n\n"
                       "# 挑戰 2：跨 deploy 檔案消失\n"
                       "FileNotFoundError: result.json\n"
                       "→ deploy 清空 workspace\n\n"
                       "# 挑戰 3：查詢語法錯誤\n"
                       "HTTP 422: Invalid character '-'\n"
                       "→ retry 3/3 → 修正語法")

    code_block(slide, R, BODY_TOP + Inches(1.72), HW, Inches(3.8), challenges_code,
               "開發過程遇到的挑戰", label_color=DANGER)


def s11_quickstart(prs):
    slide = blank(prs)
    heading(slide, "Quick Start")

    setup = ("# 1. clone\n"
             "git clone https://github.com/jason3e7/siee\n"
             "cd siee\n\n"
             "# 2. 安裝依賴\n"
             "python3 -m venv venv\n"
             "source venv/bin/activate\n"
             "pip install -r requirements.txt\n\n"
             "# 3. 設定 secret（server 端）\n"
             "export MY_API_KEY=\"sk-real-token-here\"\n\n"
             "# 4. 啟動\n"
             "python server.py      # port 5000\n"
             "python mcp_server.py  # port 5001")

    usage = ("# AI agent 部署並測試\n"
             "POST /deploy\n"
             "  files: {\"test_api.py\": \"...\"}\n\n"
             "POST /exec\n"
             "  {\"command\": \"pytest\",\n"
             "   \"args\": [\"-v\"]}\n\n"
             "GET /logs/{exec_id}\n"
             "  # 輪詢直到 status != RUNNING\n\n"
             "# 新增允許的指令\n"
             "# 編輯 server.py ALLOWED_COMMANDS")

    code_block(slide, ML, BODY_TOP + Inches(0.1), HW, Inches(4.8), setup,
               "安裝與啟動", label_color=ACCENT)
    code_block(slide, R,  BODY_TOP + Inches(0.1), HW, Inches(4.8), usage,
               "使用方式", label_color=BLUE)


def s12_summary(prs):
    slide = blank(prs)
    heading(slide, "總結")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.35),
           "SIEE 解決的核心問題：AI agent 需要在有 secret 的環境裡做真實測試，但不能讓 AI 看到 secret。",
           size=Pt(15), color=WHITE)

    cw = Inches(3.75); ch = Inches(1.8)
    cy = BODY_TOP + Inches(0.55)
    cx = ML
    for title, bc, body in [
        ("Secret 隔離",   ACCENT,
         "token 只存在 server env\n結構上就不在 AI 的 context\n無法被動洩漏"),
        ("真實整合測試",   BLUE,
         "AI 寫的 code 跑真實 API\n不用 mock，結果可信\n直接驗證整合正確性"),
        ("安全白名單",     PURPLE,
         "ALLOWED_COMMANDS 限制\n防止任意指令注入\nserver owner 完全掌控"),
    ]:
        card(slide, cx, cy, cw, ch, title=title, body=body,
             border=bc, title_color=bc, body_size=Pt(13.5))
        cx += cw + Inches(0.14)

    cy2 = cy + ch + Inches(0.2)
    cx = ML
    for title, bc, body in [
        ("MCP 原生支援",   WARN,
         "SSE transport\nClaude Code 直接接入\nAI agent 機器零額外安裝"),
        ("本機自架",       ACCENT,
         "資料不出內網\n完全掌控執行環境\n不依賴雲端服務"),
        ("23 個測試",      BLUE,
         "unit + integration 全覆蓋\n兩隻 server 獨立測試\n可信賴的 CI 基礎"),
    ]:
        card(slide, cx, cy2, cw, ch, title=title, body=body,
             border=bc, title_color=bc, body_size=Pt(13.5))
        cx += cw + Inches(0.14)


# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    s01_cover(prs)
    s02_problem(prs)
    s03_solution(prs)
    s04_philosophy(prs)
    s05_threat_model(prs)
    s06_api(prs)
    s07_allowed_commands(prs)
    s08_mcp(prs)
    s10_case_censys(prs)
    s09_case_env(prs)
    s11_quickstart(prs)
    s12_summary(prs)
    prs.save("siee.pptx")
    print(f"Saved siee.pptx  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
